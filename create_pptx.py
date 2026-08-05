"""
Generate a professional, white-background PowerPoint presentation (.pptx)
for the Optical Flow Motion Analyzer system workflow architecture.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_presentation(output_path="workflow_chart.pptx"):
    prs = Presentation()
    
    # Set slide dimensions to 16:9 widescreen (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Color Palette (White background theme with crisp modern accents)
    COLOR_BG = RGBColor(255, 255, 255)
    COLOR_CARD_BG = RGBColor(248, 250, 252) # Light slate off-white
    COLOR_BORDER = RGBColor(226, 232, 240)
    
    COLOR_TEXT_PRIMARY = RGBColor(15, 23, 42)    # Slate 900
    COLOR_TEXT_SECONDARY = RGBColor(71, 85, 105)  # Slate 600
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139)   # Slate 500
    
    # Accent Colors
    COLOR_BLUE = RGBColor(37, 99, 235)      # Lucas-Kanade / Stats
    COLOR_GREEN = RGBColor(5, 150, 105)     # Input Ingestion
    COLOR_CYAN = RGBColor(8, 145, 178)      # Preprocessing
    COLOR_PURPLE = RGBColor(124, 58, 237)   # Farneback
    COLOR_AMBER = RGBColor(217, 119, 6)     # Heatmap
    COLOR_ROSE = RGBColor(225, 29, 72)      # Object Detector

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, title_text, category_text="OPTICAL FLOW MOTION ANALYZER"):
        # Header text frame
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Category / Tag line
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_BLUE
        p_cat.font.name = "Arial"
        
        # Main Title
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_PRIMARY
        p_title.font.name = "Arial"

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    
    # Large Decorative Background Card
    rect1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    rect1.fill.solid()
    rect1.fill.fore_color.rgb = COLOR_CARD_BG
    rect1.line.color.rgb = COLOR_BORDER
    rect1.line.width = Pt(1.5)
    
    # Title Box
    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.333), Inches(3.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p0 = tf1.paragraphs[0]
    p0.text = "COMPUTER VISION PROJECT ARCHITECTURE"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_BLUE
    p0.font.name = "Arial"
    p0.space_after = Pt(14)
    
    p1 = tf1.add_paragraph()
    p1.text = "Optical Flow Motion Analyzer"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_PRIMARY
    p1.font.name = "Arial"
    p1.space_after = Pt(10)
    
    p2 = tf1.add_paragraph()
    p2.text = "System Workflow Chart, Algorithm Mechanics & Real-Time Pipeline Architecture"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_TEXT_SECONDARY
    p2.font.name = "Arial"
    p2.space_after = Pt(28)
    
    p3 = tf1.add_paragraph()
    p3.text = "TAE1 Machine Vision  |  OpenCV 4.x  |  PyQt5  |  Lucas-Kanade  |  Farnebäck  |  Heatmap"
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.font.name = "Arial"

    # =========================================================================
    # SLIDE 2: High-Level Workflow Flowchart
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "High-Level Dataflow & System Pipeline Architecture")
    
    # 5 Main Horizontal Workflow Nodes
    node_data = [
        ("1. INGESTION", "Video Source", "• Live Webcam (0)\n• MP4 / AVI File\n• QTimer Sync", COLOR_GREEN),
        ("2. PREPROCESS", "Frame Buffer", "• BGR → Gray\n• Gaussian Blur\n• Prev/Curr Buffer", COLOR_CYAN),
        ("3. ENGINE", "Optical Flow", "• Lucas-Kanade\n• Farnebäck Dense\n• Motion Heatmap", COLOR_BLUE),
        ("4. ANALYTICS", "Metrics Calculation", "• Mean/Max Speed\n• Area Coverage %\n• Polar Direction", COLOR_PURPLE),
        ("5. GUI & UI", "Canvas & Tuning", "• Qt Pixmap Render\n• Live Sliders\n• PNG Exporter", COLOR_TEXT_PRIMARY)
    ]
    
    start_x = Inches(0.8)
    node_w = Inches(2.1)
    node_h = Inches(4.5)
    gap = Inches(0.3)
    
    for i, (stage, title, details, color) in enumerate(node_data):
        x = start_x + i * (node_w + gap)
        y = Inches(1.8)
        
        # Node Card
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, node_w, node_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(2.0)
        
        # Header Badge Inside Card
        badge = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.1), y + Inches(0.1), node_w - Inches(0.2), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = stage
        p_b.font.size = Pt(11)
        p_b.font.bold = True
        p_b.font.color.rgb = RGBColor(255, 255, 255)
        p_b.alignment = PP_ALIGN.CENTER
        
        # Content Inside Card
        tb = slide2.shapes.add_textbox(x + Inches(0.15), y + Inches(0.7), node_w - Inches(0.3), node_h - Inches(0.8))
        tf_c = tb.text_frame
        tf_c.word_wrap = True
        
        p_t = tf_c.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_PRIMARY
        p_t.space_after = Pt(10)
        
        p_d = tf_c.add_paragraph()
        p_d.text = details
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_SECONDARY
        p_d.line_spacing = 1.3
        
        # Arrow connecting to next node (except last node)
        if i < len(node_data) - 1:
            arrow = slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + node_w + Inches(0.05), y + Inches(2.0), Inches(0.2), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_TEXT_MUTED
            arrow.line.fill.background()
            
    # Bottom Feedback Banner
    fb_rect = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.5), Inches(11.733), Inches(0.5))
    fb_rect.fill.solid()
    fb_rect.fill.fore_color.rgb = RGBColor(239, 246, 255)
    fb_rect.line.color.rgb = COLOR_BLUE
    fb_rect.line.width = Pt(1.0)
    
    tf_fb = fb_rect.text_frame
    p_fb = tf_fb.paragraphs[0]
    p_fb.text = "🔄 Real-time Parameter Tuning & Signal Feedback Loop: GUI Sliders ➔ Engine Parameters ➔ Live Canvas Re-render"
    p_fb.font.size = Pt(11)
    p_fb.font.bold = True
    p_fb.font.color.rgb = COLOR_BLUE
    p_fb.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 3: Stage 1 & 2 - Video Ingestion & Preprocessing
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "Stage 1 & 2: Ingestion & Preprocessing Pipeline", "CORE MODULES: VIDEO_SOURCE.PY")
    
    # Left Card - Video Source Ingestion
    card_l = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = COLOR_CARD_BG
    card_l.line.color.rgb = COLOR_GREEN
    card_l.line.width = Pt(2.0)
    
    tb_l = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "1. Video Ingestion Layer (core/video_source.py)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    p.space_after = Pt(14)
    
    bullets_l = [
        ("Dual Stream Abstraction", "Supports seamless switching between live USB webcam (device 0) and video files (.mp4, .avi, .mkv, .mov)."),
        ("OpenCV Integration", "Wraps cv2.VideoCapture to safely manage frame reading, stream closing, and error handling."),
        ("FPS Synchronized Timer", "Uses PyQt5 QTimer matching source native FPS (interval = 1000 / FPS ms) for smooth frame pacing."),
        ("Loop & Seek Control", "Automatically rewinds to frame 0 upon video completion for uninterrupted looping playback.")
    ]
    for title, desc in bullets_l:
        p1 = tf_l.add_paragraph()
        p1.text = f"• {title}: "
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_TEXT_PRIMARY
        
        # append desc
        run = p1.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(11)
        run.font.color.rgb = COLOR_TEXT_SECONDARY
        p1.space_after = Pt(8)

    # Right Card - Preprocessing Pipeline
    card_r = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.2))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLOR_CARD_BG
    card_r.line.color.rgb = COLOR_CYAN
    card_r.line.width = Pt(2.0)
    
    tb_r = slide3.shapes.add_textbox(Inches(7.133), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "2. Frame Preprocessing & Buffer"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(14)
    
    bullets_r = [
        ("Color Space Conversion", "Converts 3-channel BGR NumPy matrices into 1-channel Grayscale via cv2.cvtColor(BGR2GRAY)."),
        ("Gaussian Noise Reduction", "Applies 5x5 Gaussian Kernel blur filtering to remove sensor noise before derivative calculations."),
        ("Temporal Buffer Memory", "Stores previous frame matrix (_prev_gray) to enable consecutive frame matrix differencing."),
        ("Resolution Normalization", "Extracts image height, width, and aspect ratio for dynamic overlay bounds calculation.")
    ]
    for title, desc in bullets_r:
        p1 = tf_r.add_paragraph()
        p1.text = f"• {title}: "
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_TEXT_PRIMARY
        
        run = p1.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(11)
        run.font.color.rgb = COLOR_TEXT_SECONDARY
        p1.space_after = Pt(8)

    # =========================================================================
    # SLIDE 4: Stage 3 - Optical Flow Algorithms (3 Cards)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Stage 3: Core Optical Flow Processing Engine", "THE 3 OPTICAL FLOW ALGORITHM MODES")
    
    col_w = Inches(3.644)
    gap = Inches(0.4)
    start_x = Inches(0.8)
    
    algos = [
        ("LUCAS-KANADE", "Sparse Optical Flow", COLOR_BLUE, [
            ("Feature Detector", "Shi-Tomasi corner detection (cv2.goodFeaturesToTrack)."),
            ("Pyramidal Flow", "Tracks displacements using cv2.calcOpticalFlowPyrLK."),
            ("Rolling Trails", "Maintains up to 25 historical trajectory points per feature."),
            ("Auto Replenish", "Re-detects features when active points drop below 30.")
        ]),
        ("FARNEBÄCK", "Dense Per-Pixel Flow", COLOR_PURPLE, [
            ("Polynomial Expansion", "Computes (dx, dy) motion vector matrix for every single pixel."),
            ("Polar Conversion", "Converts vectors to magnitude and angle (cv2.cartToPolar)."),
            ("HSV Encoding", "Hue = Motion Angle, Brightness = Motion Velocity Magnitude."),
            ("Arrow Grid Overlay", "Renders direction arrows on configurable pixel grid (16px).")
        ]),
        ("HEATMAP", "Cumulative Motion Heat", COLOR_AMBER, [
            ("Frame Difference", "Calculates absolute frame intensity delta |I(t) - I(t-1)|."),
            ("Temporal Decay", "Accumulates heat with exponential decay: Acc = Acc * 0.95 + Diff."),
            ("Noise Threshold", "Suppresses background flicker below pixel threshold (15)."),
            ("False Colormap", "Maps heat to INFERNO, JET, HOT, TURBO, or VIRIDIS spectrums.")
        ])
    ]
    
    for i, (mode, name, color, details) in enumerate(algos):
        x = start_x + i * (col_w + gap)
        y = Inches(1.6)
        
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(2.0)
        
        tb = slide4.shapes.add_textbox(x + Inches(0.15), y + Inches(0.2), col_w - Inches(0.3), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = mode
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = name
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_PRIMARY
        p2.space_after = Pt(12)
        
        for item_title, item_desc in details:
            pi = tf.add_paragraph()
            pi.text = f"• {item_title}: "
            pi.font.bold = True
            pi.font.size = Pt(11)
            pi.font.color.rgb = COLOR_TEXT_PRIMARY
            
            run = pi.add_run()
            run.text = item_desc
            run.font.bold = False
            run.font.size = Pt(10)
            run.font.color.rgb = COLOR_TEXT_SECONDARY
            pi.space_after = Pt(6)

    # =========================================================================
    # SLIDE 5: Auxiliary Stage - MOG2 Object Detector
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Auxiliary Subsystem: MOG2 Motion Object Detector", "PARALLEL ANALYSIS: CORE/MOTION_DETECTOR.PY")
    
    # Large Card
    card_det = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    card_det.fill.solid()
    card_det.fill.fore_color.rgb = COLOR_CARD_BG
    card_det.line.color.rgb = COLOR_ROSE
    card_det.line.width = Pt(2.0)
    
    tb_det = slide5.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.133), Inches(4.8))
    tf_det = tb_det.text_frame
    tf_det.word_wrap = True
    
    p = tf_det.paragraphs[0]
    p.text = "Background Subtraction & Object Centroid Tracking Pipeline"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROSE
    p.space_after = Pt(16)
    
    det_steps = [
        ("MOG2 Background Subtraction", "Uses Mixture of Gaussians (cv2.createBackgroundSubtractorMOG2) with history=300 and shadow suppression (shadow pixels set to 0)."),
        ("Morphological Mask Cleanup", "Executes Morphological Opening (remove noise), Closing (fill holes), and Dilation (join fragmented contours) using 5x5 ellipse kernel."),
        ("Contour Detection & Filtering", "Extracts object contours (cv2.findContours) and filters out noise blobs smaller than min_area (default 800 px²)."),
        ("Centroid Speed & Direction", "Calculates centroid displacements between frames to measure object speed (px/frame) and 360° motion angle."),
        ("Cardinal Arrow Annotations", "Maps motion angle to 8 cardinal directions (↗, ↑, ↖, ←, ↙, ↓, ↘, →) and renders HUD bounding boxes with motion trails.")
    ]
    
    for step_title, step_desc in det_steps:
        p_step = tf_det.add_paragraph()
        p_step.text = f"• {step_title}: "
        p_step.font.bold = True
        p_step.font.size = Pt(13)
        p_step.font.color.rgb = COLOR_TEXT_PRIMARY
        
        run = p_step.add_run()
        run.text = step_desc
        run.font.bold = False
        run.font.size = Pt(12)
        run.font.color.rgb = COLOR_TEXT_SECONDARY
        p_step.space_after = Pt(10)

    # =========================================================================
    # SLIDE 6: Stage 4 & 5 - Analytics & GUI Presentation
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "Stage 4 & 5: Analytics Engine & PyQt5 GUI Presentation", "ANALYTICS & USER INTERFACE LAYER")
    
    # Left Box - Analytics
    card_an = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    card_an.fill.solid()
    card_an.fill.fore_color.rgb = COLOR_CARD_BG
    card_an.line.color.rgb = COLOR_BLUE
    card_an.line.width = Pt(2.0)
    
    tb_an = slide6.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_an = tb_an.text_frame
    tf_an.word_wrap = True
    
    p = tf_an.paragraphs[0]
    p.text = "4. Motion Analytics (motion_stats.py)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    p.space_after = Pt(14)
    
    an_items = [
        ("Average Speed Magnitude", "Calculates arithmetic mean velocity across all active motion vectors."),
        ("Maximum Speed", "Tracks peak velocity magnitude in current frame."),
        ("Active Area Percentage", "Percentage of total frame area containing active motion (magnitude > threshold)."),
        ("8-Bin Polar Histogram", "Divides 360° into 8 octants to compute angular distribution of flow vectors.")
    ]
    for title, desc in an_items:
        p1 = tf_an.add_paragraph()
        p1.text = f"• {title}: "
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_TEXT_PRIMARY
        
        run = p1.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(11)
        run.font.color.rgb = COLOR_TEXT_SECONDARY
        p1.space_after = Pt(8)

    # Right Box - GUI Canvas
    card_gui = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.2))
    card_gui.fill.solid()
    card_gui.fill.fore_color.rgb = COLOR_CARD_BG
    card_gui.line.color.rgb = COLOR_TEXT_PRIMARY
    card_gui.line.width = Pt(2.0)
    
    tb_gui = slide6.shapes.add_textbox(Inches(7.133), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_gui = tb_gui.text_frame
    tf_gui.word_wrap = True
    
    p = tf_gui.paragraphs[0]
    p.text = "5. PyQt5 GUI & Canvas (gui/)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_PRIMARY
    p.space_after = Pt(14)
    
    gui_items = [
        ("BGR to Qt Pixmap Conversion", "Converts OpenCV NumPy matrices into QImage & QPixmap for hardware-accelerated rendering."),
        ("Video Canvas Display", "Renders scaled output frames with top-left HUD overlay (FPS, Mode, Object Count)."),
        ("Live Parameter Tuning", "PyQt Sliders bind real-time parameters (LK window size, Farnebäck scale, Heatmap decay)."),
        ("Screenshot & Export", "One-click save action captures current processed frame directly to PNG image files.")
    ]
    for title, desc in gui_items:
        p1 = tf_gui.add_paragraph()
        p1.text = f"• {title}: "
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_TEXT_PRIMARY
        
        run = p1.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(11)
        run.font.color.rgb = COLOR_TEXT_SECONDARY
        p1.space_after = Pt(8)

    # =========================================================================
    # SLIDE 7: Parameter Signal Routing Table
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Parameter Signal Routing & GUI Slider Mapping", "TECHNICAL REFERENCE SPECIFICATION")
    
    # Table Shape
    rows, cols = 8, 5
    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(11.733)
    height = Inches(5.2)
    
    table_shape = slide7.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.4)
    table.columns[3].width = Inches(2.0)
    table.columns[4].width = Inches(3.133)
    
    table_data = [
        ["Module", "GUI Widget / Slider", "Target Component", "Parameter", "Algorithmic Impact"],
        ["Lucas-Kanade", "Window Size Slider", "LucasKanadeTracker", "winSize (N, N)", "Neighborhood integration window size"],
        ["Lucas-Kanade", "Max Corners Slider", "LucasKanadeTracker", "maxCorners", "Cap on Shi-Tomasi feature points count"],
        ["Farnebäck", "Pyramid Scale Slider", "FarnebackAnalyzer", "pyr_scale", "Multi-resolution pyramid scale factor"],
        ["Farnebäck", "Arrow Step Grid", "FarnebackAnalyzer", "arrow_step", "Pixel grid interval between vector arrows"],
        ["Heatmap", "Decay Rate Slider", "MotionHeatmap", "decay (0.80–0.99)", "Persistence decay speed of motion trails"],
        ["Heatmap", "Colormap Combo Box", "MotionHeatmap", "colormap_name", "False-color spectrum (INFERNO, JET, etc.)"],
        ["Object Detector", "Min Contour Area", "MotionDetector", "min_area", "Noise suppression threshold (default 800px²)"]
    ]
    
    for r_idx, row in enumerate(table_data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TEXT_PRIMARY
                p.font.bold = True
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(255, 255, 255)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_CARD_BG if r_idx % 2 == 0 else RGBColor(255, 255, 255)
                p.font.size = Pt(10)
                p.font.color.rgb = COLOR_TEXT_PRIMARY

    # Save presentation
    prs.save(output_path)
    print(f"Successfully generated PowerPoint presentation: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
